from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x12, 0x0A, 0x02)   # deep amber-dark
AMBER       = RGBColor(0xB0, 0x7A, 0x10)   # accent gold
AMBER_LIGHT = RGBColor(0xFF, 0xEC, 0xB2)   # tile colour
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xFF, 0xF4, 0xE0)
SLATE       = RGBColor(0x3C, 0x2A, 0x08)
MID_AMBER   = RGBColor(0xD4, 0xA0, 0x30)
LIGHT_GREY  = RGBColor(0xF5, 0xF0, 0xE8)
CARD_BG     = RGBColor(0xFF, 0xF0, 0xD0)

def blank_slide():
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox

def add_multiline(slide, lines, l, t, w, h, size, color=WHITE, bold_first=False, spacing=1.15, font="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
        run.font.name = font
        p.space_after = Pt(size * 0.3)
    return txBox

def accent_bar(slide, y=0.55):
    add_rect(slide, 0, y, 13.33, 0.04, AMBER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s1 = blank_slide()
fill_bg(s1, DARK_BG)

# full-width gradient-like band at top
add_rect(s1, 0, 0, 13.33, 2.6, RGBColor(0x1E, 0x13, 0x05))
accent_bar(s1, 0.0)

# big logo mark
add_text(s1, "✦", 0.7, 0.3, 1.2, 1.1, 52, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# App name
add_text(s1, "SEFER", 1.7, 0.22, 6, 1.0, 58, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="Georgia")
add_text(s1, "A Sacred Space for Scripture, Reflection & Prayer", 1.7, 1.1, 9.5, 0.7, 17, bold=False,
         color=MID_AMBER, align=PP_ALIGN.LEFT, font="Calibri")

accent_bar(s1, 2.52)

# tag line block
add_text(s1, "Built from scratch. Designed for devotion.", 0.7, 2.85, 10, 0.7, 22,
         bold=True, color=OFF_WHITE, align=PP_ALIGN.LEFT, font="Georgia")

add_text(s1, "A premium desktop Bible study app combining the full 31,000-verse Scripture database\n"
             "with rich annotation, prayer journalling, cross-references, and beautiful themes — "
             "all stored securely on your Mac with iCloud sync.",
         0.7, 3.55, 11.5, 1.5, 13.5, color=RGBColor(0xCC, 0xB4, 0x80), align=PP_ALIGN.LEFT)

# stats strip at bottom
add_rect(s1, 0, 5.55, 13.33, 1.55, RGBColor(0x0C, 0x07, 0x01))
stats = [("31,000+", "Bible Verses"), ("10", "Themes"), ("3", "Core Modules"),
         ("100%", "Private & Offline"), ("iCloud", "Auto-Sync")]
for i, (val, lbl) in enumerate(stats):
    x = 0.5 + i * 2.6
    add_text(s1, val, x, 5.7, 2.3, 0.6, 22, bold=True, color=AMBER, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s1, lbl, x, 6.25, 2.3, 0.5, 10, color=RGBColor(0xAA, 0x90, 0x50), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════════════════
s2 = blank_slide()
fill_bg(s2, DARK_BG)
add_rect(s2, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s2, 0.72)

add_text(s2, "The Problem with Today's Bible Apps", 0.6, 0.1, 12, 0.65, 24,
         bold=True, color=WHITE, font="Georgia")
add_text(s2, "Most apps force you to choose between beauty, depth, or privacy — Sefer refuses to compromise.",
         0.6, 0.78, 12, 0.45, 12, color=MID_AMBER)

problems = [
    ("📱  Web & Mobile First",
     "Popular apps like YouVersion and Logos are designed for phones or browsers.\nDesktop users are second-class citizens with slow, bloated interfaces."),
    ("🔒  Your Data Isn't Yours",
     "Most apps sync notes and highlights to their own servers.\nYour personal reflections, prayer points, and annotations are stored by a company you don't control."),
    ("💡  Shallow Study Tools",
     "Basic apps offer reading plans but lack deep cross-reference networks,\nrich text annotation, or integrated prayer journalling."),
    ("🎨  One-Size-Fits-All Design",
     "A single default theme with no personalisation.\nNo consideration for how different users prefer to read and reflect."),
    ("🌐  Always-Online Dependency",
     "Many apps require an internet connection for even basic reading.\nNo offline-first architecture."),
]
for i, (title, body) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.38 + row * 1.72
    add_rect(s2, x, y, 6.1, 1.52, RGBColor(0x22, 0x16, 0x05), RGBColor(0x6B, 0x4A, 0x10), 0.8)
    add_text(s2, title, x + 0.18, y + 0.12, 5.7, 0.42, 12, bold=True, color=AMBER)
    add_text(s2, body, x + 0.18, y + 0.52, 5.7, 0.88, 10.5, color=RGBColor(0xCC, 0xB8, 0x88))

# last item centred
add_rect(s2, 4.1, 6.82 - 1.72, 5.1, 1.52, RGBColor(0x22, 0x16, 0x05), RGBColor(0x6B, 0x4A, 0x10), 0.8)
# already placed as row 2 — adjust: skip, 5th item goes at row=2 col=0 → (0.5, 5.08)
# Re-draw 5th problem
px, py = 3.1, 4.9
add_rect(s2, px, py, 7.1, 1.52, RGBColor(0x22, 0x16, 0x05), RGBColor(0x6B, 0x4A, 0x10), 0.8)
add_text(s2, problems[4][0], px+0.18, py+0.12, 6.7, 0.42, 12, bold=True, color=AMBER)
add_text(s2, problems[4][1], px+0.18, py+0.52, 6.7, 0.88, 10.5, color=RGBColor(0xCC, 0xB8, 0x88))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introducing Sefer
# ════════════════════════════════════════════════════════════════════════════
s3 = blank_slide()
fill_bg(s3, DARK_BG)
add_rect(s3, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s3, 0.72)

add_text(s3, "Introducing Sefer", 0.6, 0.1, 10, 0.65, 28, bold=True, color=WHITE, font="Georgia")
add_text(s3, "Hebrew for 'Book' — a Mac-native Bible study app built for depth, beauty, and privacy.",
         0.6, 0.78, 12, 0.45, 13, color=MID_AMBER)

add_text(s3, "\"Sefer\" (סֵפֶר) means Book in Hebrew — the same word used for the Torah scrolls.\n"
             "We named the app after the holiest of objects to signal our intent: this is not a casual reading app.\n"
             "It is a study companion built for people who take Scripture seriously.",
         0.7, 1.32, 11.8, 1.1, 13, color=OFF_WHITE,
         italic=True)

pillars = [
    ("📖", "Complete Scripture", "Full 31,000-verse Bible with NKJV and multiple translation support. Every word, every book, instantly accessible."),
    ("🔍", "Deep Study Tools", "31,000-verse cross-reference network, verse highlights in 5 colours, rich-text notes with fonts and formatting."),
    ("🙏", "Prayer & Proclamation", "Integrated prayer journal with personal notes and a daily Proclamations panel for faith declarations."),
    ("🎨", "10 Beautiful Themes", "Shiloh, Horeb, Galilee, Edom, Sharon, Moriah, Sinai, Carmel, Patmos, Tyre — each with a biblical reference."),
    ("🔒", "Private by Design", "Zero cloud accounts required. All data lives on your Mac with optional iCloud Drive backup you control."),
    ("✦", "Native Mac App", "Built with Electron for a true desktop experience. Single-instance, offline-first, fast, and always available."),
]
for i, (icon, title, body) in enumerate(pillars):
    col = i % 3
    row = i // 3
    x = 0.45 + col * 4.26
    y = 2.6 + row * 2.25
    add_rect(s3, x, y, 3.9, 2.0, RGBColor(0x1E, 0x14, 0x04), RGBColor(0x8A, 0x62, 0x18), 0.6)
    add_text(s3, icon, x + 0.18, y + 0.18, 0.6, 0.5, 22, align=PP_ALIGN.CENTER)
    add_text(s3, title, x + 0.78, y + 0.15, 2.9, 0.42, 12, bold=True, color=AMBER)
    add_text(s3, body, x + 0.18, y + 0.65, 3.55, 1.15, 10, color=RGBColor(0xCC, 0xB4, 0x7A))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Bible Reader
# ════════════════════════════════════════════════════════════════════════════
s4 = blank_slide()
fill_bg(s4, DARK_BG)
add_rect(s4, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s4, 0.72)

add_text(s4, "Module 1 — Bible Reader", 0.6, 0.1, 12, 0.65, 26, bold=True, color=WHITE, font="Georgia")
add_text(s4, "The heart of Sefer — a distraction-free reading environment built for serious study.",
         0.6, 0.78, 12, 0.42, 12.5, color=MID_AMBER)

# Left panel - features list
features = [
    ("Full 39 + 27 Book Library", "Every Old and New Testament book, all 31,000 verses, loaded locally. No internet required."),
    ("Multiple Translations", "Switch between NKJV, KJV, ESV, NIV, NLT and more with a single click per chapter."),
    ("Verse-Level Actions", "Tap any verse number to highlight (5 colours), write a note, view cross-references, or compare translations."),
    ("Smart Search", "Search by keyword, emotion (anxious, hopeful, loved…), or direct reference like 'John 3:16'."),
    ("Cross-Reference Popup", "31,000-verse cross-reference network — click any verse to instantly see related passages with verse previews."),
    ("Chapter Navigation", "Prev/Next chapter buttons, chapter jump dropdown, book breadcrumb — always one click from anywhere."),
    ("Reading Hints", "Verse numbers display inline; a subtle footer reminds you how to interact on first open."),
]
for i, (title, body) in enumerate(features):
    y = 1.32 + i * 0.82
    add_rect(s4, 0.5, y, 0.06, 0.42, AMBER)
    add_text(s4, title, 0.72, y, 5.5, 0.32, 11, bold=True, color=AMBER_LIGHT)
    add_text(s4, body, 0.72, y + 0.3, 5.5, 0.42, 10, color=RGBColor(0xBB, 0xA0, 0x68))

# Right panel - mock reader UI
add_rect(s4, 7.1, 1.25, 5.8, 5.9, RGBColor(0xFF, 0xF0, 0xD0), RGBColor(0xB0, 0x7A, 0x10), 1.2)
add_rect(s4, 7.1, 1.25, 5.8, 0.52, RGBColor(0xFF, 0xE8, 0xB8))
add_text(s4, "← Genesis          NKJV ▾", 7.25, 1.29, 5.5, 0.38, 9.5, color=RGBColor(0x6B, 0x4A, 0x10))
add_text(s4, "Genesis 1", 7.3, 1.9, 5.5, 0.55, 20, bold=True, color=RGBColor(0x6B, 0x48, 0x00), font="Georgia")
mock_verses = [
    ("1", "In the beginning God created the heavens and the earth."),
    ("2", "The earth was without form, and void; and darkness was on the face of the deep."),
    ("3", 'Then God said, “Let there be light”; and there was light.'),
    ("4", 'And God saw the light, that it was good; and God divided the light from the darkness.'),
    ("5", "God called the light Day, and the darkness He called Night."),
]
for i, (vnum, vtxt) in enumerate(mock_verses):
    y = 2.55 + i * 0.72
    add_rect(s4, 7.22, y + 0.04, 0.32, 0.26, AMBER if vnum=="3" else RGBColor(0xE8, 0xD0, 0x98))
    add_text(s4, vnum, 7.22, y+0.02, 0.32, 0.28, 7.5, bold=True, color=DARK_BG if vnum=="3" else RGBColor(0x8A,0x62,0x18), align=PP_ALIGN.CENTER)
    wrap = vtxt[:72] + ("…" if len(vtxt) > 72 else "")
    add_text(s4, wrap, 7.6, y, 5.05, 0.6, 9.5, color=RGBColor(0x18, 0x0E, 0x02))

add_text(s4, "↑ Tap verse number to highlight, annotate or cross-reference", 7.22, 6.85, 5.6, 0.3,
         8, italic=True, color=RGBColor(0x8A, 0x68, 0x28), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Cross-Reference Engine
# ════════════════════════════════════════════════════════════════════════════
s5 = blank_slide()
fill_bg(s5, DARK_BG)
add_rect(s5, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s5, 0.72)

add_text(s5, "Cross-Reference Engine", 0.6, 0.1, 12, 0.65, 26, bold=True, color=WHITE, font="Georgia")
add_text(s5, "31,000 curated theological connections — the largest cross-reference network in any desktop Bible app.",
         0.6, 0.78, 12.5, 0.42, 12.5, color=MID_AMBER)

# big number
add_rect(s5, 0.5, 1.35, 5.8, 3.4, RGBColor(0x1E, 0x14, 0x04), RGBColor(0x8A, 0x62, 0x18), 0.8)
add_text(s5, "31,000+", 0.6, 1.5, 5.6, 1.2, 52, bold=True, color=AMBER, align=PP_ALIGN.CENTER, font="Georgia")
add_text(s5, "Cross-Reference Pairs", 0.6, 2.65, 5.6, 0.5, 14, color=OFF_WHITE, align=PP_ALIGN.CENTER)
add_text(s5, "Spanning Law → Prophets → Gospels → Epistles\n"
             "Covering every major theological theme:\nRedemption · Light · Covenant · Faith · Praise",
         0.7, 3.2, 5.5, 1.4, 11, color=RGBColor(0xCC, 0xB4, 0x7A), align=PP_ALIGN.CENTER)

# how it works
add_text(s5, "How It Works", 6.8, 1.35, 6.3, 0.4, 14, bold=True, color=AMBER_LIGHT)
steps = [
    ("1", "Tap any verse number in the reader"),
    ("2", "Select 'Cross-references' from the action menu"),
    ("3", "A modal pops up instantly — no loading, no internet"),
    ("4", "Each reference shows book, chapter:verse + live snippet"),
    ("5", "Click any reference to navigate there directly"),
]
for i, (n, txt) in enumerate(steps):
    y = 1.9 + i * 0.76
    add_rect(s5, 6.8, y, 0.38, 0.38, AMBER)
    add_text(s5, n, 6.8, y, 0.38, 0.38, 13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s5, txt, 7.28, y + 0.04, 5.7, 0.34, 11.5, color=OFF_WHITE)

# example popup mock
add_rect(s5, 6.8, 5.65, 6.2, 1.65, RGBColor(0xFF, 0xF0, 0xD0), RGBColor(0xB0, 0x7A, 0x10), 1)
add_text(s5, "CROSS-REFERENCES  Genesis 1:3", 7.0, 5.75, 5.8, 0.3, 9, bold=True, color=RGBColor(0x6B, 0x48, 0x00))
refs = [("John 1:5", '"The light shines in darkness…"'), ("2 Cor 4:6", '"God commanded light to shine…"'), ("Isaiah 60:1", '"Arise, shine; your light has come!"')]
for i, (ref, snip) in enumerate(refs):
    y = 6.1 + i * 0.37
    add_text(s5, ref, 7.0, y, 1.5, 0.32, 9, bold=True, color=AMBER)
    add_text(s5, snip, 8.55, y, 4.2, 0.32, 8.5, italic=True, color=RGBColor(0x55, 0x3A, 0x10))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Notes System
# ════════════════════════════════════════════════════════════════════════════
s6 = blank_slide()
fill_bg(s6, DARK_BG)
add_rect(s6, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s6, 0.72)

add_text(s6, "Module 2 — Notes System", 0.6, 0.1, 12, 0.65, 26, bold=True, color=WHITE, font="Georgia")
add_text(s6, "Two types of notes. One rich, unified journal. Export to Word or plain text anytime.",
         0.6, 0.78, 12, 0.42, 12.5, color=MID_AMBER)

# Personal notes col
add_rect(s6, 0.5, 1.32, 5.9, 0.48, RGBColor(0x2A, 0x1A, 0x06))
add_text(s6, "📓  Personal Notes", 0.65, 1.36, 5.6, 0.38, 13, bold=True, color=AMBER_LIGHT)

personal_pts = [
    "Free-form notes on any topic — sermons, devotions, reflections",
    "Title + rich text body with Bold, Italic, Underline formatting",
    "9 font families (Lora, Cormorant, Playfair, Georgia…)",
    "Font sizes 10–24pt in steps of 2",
    "Bullet lists and numbered lists",
    "7 colour labels for personal organisation",
    "Uniform card grid — equal-height tiles for clean layout",
    "Export each note as .docx or .txt",
]
for i, pt in enumerate(personal_pts):
    add_text(s6, f"  ◆  {pt}", 0.5, 1.88 + i * 0.42, 5.9, 0.38, 10.5, color=RGBColor(0xCC, 0xB4, 0x7A))

# Verse notes col
add_rect(s6, 7.1, 1.32, 5.9, 0.48, RGBColor(0x2A, 0x1A, 0x06))
add_text(s6, "📖  Verse Notes", 7.25, 1.36, 5.6, 0.38, 13, bold=True, color=AMBER_LIGHT)

verse_pts = [
    "Linked directly to a Bible verse (e.g. John 3:16)",
    "Verse text shown as a styled quote in the card header",
    "Multiple notes per verse — each with timestamp",
    "Created directly from the Bible reader's action menu",
    "📖 button to jump back to that verse in the reader",
    "Colour-coded accent inherited from note label",
    "Full-text search across all verse notes",
    "Export as structured Word document",
]
for i, pt in enumerate(verse_pts):
    add_text(s6, f"  ◆  {pt}", 7.1, 1.88 + i * 0.42, 5.9, 0.38, 10.5, color=RGBColor(0xCC, 0xB4, 0x7A))

# divider
add_rect(s6, 6.56, 1.32, 0.06, 5.8, AMBER)

# bottom strip
add_rect(s6, 0, 6.8, 13.33, 0.7, RGBColor(0x0C, 0x07, 0x01))
add_text(s6, "All notes are stored locally in JSON — no account, no server, no subscription required.",
         0.6, 6.87, 12.2, 0.4, 11.5, italic=True, color=MID_AMBER, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Prayer & Proclamations
# ════════════════════════════════════════════════════════════════════════════
s7 = blank_slide()
fill_bg(s7, DARK_BG)
add_rect(s7, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s7, 0.72)

add_text(s7, "Module 3 — Prayer & Proclamations", 0.6, 0.1, 12, 0.65, 24, bold=True, color=WHITE, font="Georgia")
add_text(s7, "A dedicated space to bring your requests before God and declare His word over your life.",
         0.6, 0.78, 12, 0.42, 12.5, color=MID_AMBER)

# Prayer journal block
add_rect(s7, 0.5, 1.32, 5.9, 5.85, RGBColor(0x18, 0x10, 0x04))
add_rect(s7, 0.5, 1.32, 5.9, 0.5, RGBColor(0x2A, 0x1A, 0x06))
add_text(s7, "🙏  Prayer Journal", 0.65, 1.35, 5.6, 0.4, 13, bold=True, color=AMBER_LIGHT)

prayer_pts = [
    ("Create prayer requests", "Give each request a title and a rich-text body with full formatting support."),
    ("Colour labels", "7 colours to categorise — family, health, ministry, work, etc."),
    ("Personal context", "Add your name and the person you're praying for on each card."),
    ("Date tracking", "Every prayer shows when it was written and last updated."),
    ("Edit & delete", "Full lifecycle management — nothing is permanent unless you want it to be."),
    ("Persistent storage", "Saved to localStorage + local file + iCloud Drive simultaneously."),
]
for i, (title, body) in enumerate(prayer_pts):
    y = 1.95 + i * 0.85
    add_rect(s7, 0.6, y, 0.06, 0.5, AMBER)
    add_text(s7, title, 0.8, y, 5.4, 0.32, 11, bold=True, color=AMBER_LIGHT)
    add_text(s7, body, 0.8, y + 0.3, 5.4, 0.44, 10, color=RGBColor(0xBB, 0xA0, 0x68))

# Proclamations block
add_rect(s7, 7.1, 1.32, 5.9, 5.85, RGBColor(0x18, 0x10, 0x04))
add_rect(s7, 7.1, 1.32, 5.9, 0.5, RGBColor(0x2A, 0x1A, 0x06))
add_text(s7, "📣  Proclamations Panel", 7.25, 1.35, 5.6, 0.4, 13, bold=True, color=AMBER_LIGHT)

proc_pts = [
    ("Add any Bible verse", "Type a reference (e.g. Philippians 4:13) and Sefer fetches and displays the full text."),
    ("Declare daily", "Read each card aloud as a faith declaration — a discipline rooted in biblical tradition."),
    ("Live verse lookup", "Verse text is retrieved and cached automatically — no typing required."),
    ("Open in reader", "📖 button takes you directly to that verse in the Bible reader."),
    ("Persistent & synced", "Your proclamation list is saved locally and backed up to iCloud Drive."),
    ("Clean card layout", "Each card shows reference + full verse text in a legible, distraction-free design."),
]
for i, (title, body) in enumerate(proc_pts):
    y = 1.95 + i * 0.85
    add_rect(s7, 7.2, y, 0.06, 0.5, AMBER)
    add_text(s7, title, 7.4, y, 5.4, 0.32, 11, bold=True, color=AMBER_LIGHT)
    add_text(s7, body, 7.4, y + 0.3, 5.4, 0.44, 10, color=RGBColor(0xBB, 0xA0, 0x68))

add_rect(s7, 6.56, 1.32, 0.06, 5.85, AMBER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Themes & Design
# ════════════════════════════════════════════════════════════════════════════
s8 = blank_slide()
fill_bg(s8, DARK_BG)
add_rect(s8, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s8, 0.72)

add_text(s8, "10 Biblically-Named Themes", 0.6, 0.1, 12, 0.65, 26, bold=True, color=WHITE, font="Georgia")
add_text(s8, "Every theme is named after a place in Scripture — each with its own palette, contrast-tested for readability.",
         0.6, 0.78, 12.5, 0.42, 12.5, color=MID_AMBER)

themes = [
    ("Shiloh",   RGBColor(0xB0,0x7A,0x10), RGBColor(0xFF,0xF1,0xC4), "Deep amber honey"),
    ("Horeb",    RGBColor(0x68,0x78,0xA8), RGBColor(0xE0,0xE6,0xF6), "Graphite storm blue"),
    ("Galilee",  RGBColor(0x1A,0x5E,0xC8), RGBColor(0xC8,0xE2,0xFF), "Midnight ocean"),
    ("Edom",     RGBColor(0xC0,0x18,0x18), RGBColor(0xFF,0xD9,0xD9), "Deep crimson wine"),
    ("Sharon",   RGBColor(0xA8,0x28,0x60), RGBColor(0xFF,0xD5,0xEA), "Berry magenta rose"),
    ("Moriah",   RGBColor(0x50,0x48,0xC8), RGBColor(0xE0,0xDC,0xFF), "Deep indigo"),
    ("Sinai",    RGBColor(0x5A,0x78,0x20), RGBColor(0xE2,0xF4,0xC9), "Sage desert green"),
    ("Carmel",   RGBColor(0x1A,0x70,0x38), RGBColor(0xC9,0xF5,0xDC), "Emerald forest"),
    ("Patmos",   RGBColor(0x6C,0x28,0xC0), RGBColor(0xE4,0xD2,0xFF), "Deep amethyst"),
    ("Tyre",     RGBColor(0x10,0x78,0x78), RGBColor(0xC2,0xF5,0xF4), "Teal lagoon"),
]
for i, (name, accent, tile, desc) in enumerate(themes):
    col = i % 5
    row = i // 5
    x = 0.4 + col * 2.5
    y = 1.32 + row * 2.72
    # swatch
    add_rect(s8, x, y, 2.2, 1.35, tile, border_color=accent, border_pt=1.5)
    add_rect(s8, x, y, 2.2, 0.42, accent)
    add_text(s8, name, x + 0.1, y + 0.06, 2.0, 0.32, 12, bold=True, color=WHITE, font="Georgia")
    add_text(s8, desc, x + 0.1, y + 0.52, 2.0, 0.68, 9, color=RGBColor(0x33, 0x22, 0x0A))
    # verse reference
    verses = ["Gen 49:10","1 Kgs 19:8","Matt 4:18","Isa 63:1","Song 2:1",
              "Gen 22:2","Ex 19:18","Song 7:5","Rev 1:9","Ezek 27:3"]
    add_text(s8, verses[i], x + 0.1, y + 1.08, 2.0, 0.3, 8, color=accent, italic=True)

add_rect(s8, 0, 6.76, 13.33, 0.74, RGBColor(0x0C, 0x07, 0x01))
add_text(s8, "All themes pass WCAG AA contrast checks on both dark backgrounds and light tile surfaces.\n"
             "CSS variable system allows instant theme switching with zero page reload.",
         0.6, 6.82, 12.5, 0.55, 10.5, color=MID_AMBER, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Data & Sync
# ════════════════════════════════════════════════════════════════════════════
s9 = blank_slide()
fill_bg(s9, DARK_BG)
add_rect(s9, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s9, 0.72)

add_text(s9, "Your Data. Your Device. Your iCloud.", 0.6, 0.1, 12, 0.65, 26, bold=True, color=WHITE, font="Georgia")
add_text(s9, "Three independent layers of data protection — no account, no subscription, no server dependency.",
         0.6, 0.78, 12.5, 0.42, 12.5, color=MID_AMBER)

layers = [
    ("Layer 1", "localStorage", "Instant in-app storage.\nEvery keystroke saves immediately to the browser's\nlocalStorage inside Electron — zero latency.",
     RGBColor(0x18,0x12,0x04), RGBColor(0xB0,0x7A,0x10)),
    ("Layer 2", "Local File Backup", "Crash-proof persistence.\nEvery write triggers an atomic backup to:\n~/Library/Application Support/Sefer/sefer-storage-backup.json\nSurvives force-quit, crash, or localStorage wipe.",
     RGBColor(0x14,0x10,0x04), RGBColor(0xB0,0x7A,0x10)),
    ("Layer 3", "iCloud Drive", "Cross-device & cloud safety.\nSimultaneously writes to:\n~/Library/Mobile Documents/com~apple~CloudDocs/Sefer/\nOn startup, Sefer reads whichever copy is newer.",
     RGBColor(0x10,0x0E,0x04), RGBColor(0xB0,0x7A,0x10)),
]
for i, (label, title, body, bg, border) in enumerate(layers):
    x = 0.5 + i * 4.22
    add_rect(s9, x, 1.38, 3.9, 4.35, bg, border, 1.0)
    add_rect(s9, x, 1.38, 3.9, 0.48, RGBColor(0x2A,0x1A,0x06))
    add_text(s9, label, x+0.15, 1.41, 3.6, 0.3, 9, color=MID_AMBER, bold=True)
    add_text(s9, title, x+0.15, 1.68, 3.6, 0.42, 13, color=AMBER_LIGHT, bold=True)
    add_text(s9, body, x+0.15, 2.22, 3.6, 2.8, 10.5, color=RGBColor(0xCC,0xB4,0x7A))

# arrows between layers
for ax in [4.42, 8.64]:
    add_text(s9, "→", ax, 3.0, 0.55, 0.55, 24, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Bottom feature list
add_rect(s9, 0.5, 5.88, 12.33, 1.38, RGBColor(0x1A,0x12,0x04), RGBColor(0x6B,0x4A,0x10), 0.6)
add_text(s9, "Additional Data Features", 0.7, 5.95, 12, 0.35, 11, bold=True, color=AMBER_LIGHT)
extras = [
    "☁  Sync badge in header shows 'Synced' or 'Pending' with last-sync date on hover",
    "🔒  Single-instance lock — clicking the icon when app is open brings it to front, never duplicates",
    "📤  Export all notes as .docx (Word) or .txt with one click",
    "🔄  On startup, app restores from the newest available backup automatically",
]
for i, e in enumerate(extras):
    col, row = i % 2, i // 2
    add_text(s9, e, 0.7 + col * 6.1, 6.38 + row * 0.38, 5.9, 0.36, 10, color=RGBColor(0xCC,0xB4,0x7A))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Competitive Comparison
# ════════════════════════════════════════════════════════════════════════════
s10 = blank_slide()
fill_bg(s10, DARK_BG)
add_rect(s10, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s10, 0.72)

add_text(s10, "How Sefer Stands Out", 0.6, 0.1, 12, 0.65, 26, bold=True, color=WHITE, font="Georgia")
add_text(s10, "A feature-by-feature comparison against the leading Bible study platforms.",
         0.6, 0.78, 12.5, 0.42, 12.5, color=MID_AMBER)

# Table header
headers = ["Feature", "Sefer", "YouVersion", "Logos", "Bible Gateway", "OliveTree"]
col_widths = [3.2, 1.6, 1.6, 1.6, 1.6, 1.6]
col_starts = [0.4]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w + 0.06)

for ci, (hdr, cw, cx) in enumerate(zip(headers, col_widths, col_starts)):
    bg = AMBER if ci == 1 else (RGBColor(0x2A,0x1A,0x06) if ci == 0 else RGBColor(0x1E,0x14,0x04))
    add_rect(s10, cx, 1.32, cw, 0.44, bg)
    add_text(s10, hdr, cx+0.08, 1.35, cw-0.1, 0.38, 10.5, bold=True,
             color=DARK_BG if ci==1 else AMBER_LIGHT, align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)

rows = [
    ("Offline-first (no internet needed)",   ["✅","❌","⚠️","❌","⚠️"]),
    ("No account / subscription required",   ["✅","❌","❌","❌","❌"]),
    ("Native Mac desktop app",               ["✅","❌","✅","❌","✅"]),
    ("31,000-verse cross-reference network", ["✅","❌","✅","❌","⚠️"]),
    ("Rich-text notes with font control",    ["✅","⚠️","✅","❌","⚠️"]),
    ("Integrated prayer journal",            ["✅","✅","❌","❌","❌"]),
    ("Proclamations / declarations panel",   ["✅","❌","❌","❌","❌"]),
    ("10 custom themes",                     ["✅","❌","❌","❌","❌"]),
    ("iCloud Drive sync (no 3rd party)",     ["✅","❌","❌","❌","❌"]),
    ("Emotion-based verse search",           ["✅","❌","❌","❌","❌"]),
]
CHECK_COLOR  = RGBColor(0x2A, 0x9A, 0x40)
CROSS_COLOR  = RGBColor(0xC0, 0x30, 0x20)
WARN_COLOR   = RGBColor(0xB8, 0x88, 0x10)

for ri, (feature, vals) in enumerate(rows):
    y = 1.82 + ri * 0.48
    row_bg = RGBColor(0x18,0x11,0x04) if ri % 2 == 0 else RGBColor(0x14,0x0E,0x03)
    add_rect(s10, col_starts[0], y, col_widths[0], 0.43, row_bg)
    add_text(s10, feature, col_starts[0]+0.12, y+0.07, col_widths[0]-0.15, 0.34, 9.5, color=OFF_WHITE)
    for ci, (val, cw, cx) in enumerate(zip(vals, col_widths[1:], col_starts[1:])):
        bg = RGBColor(0x22,0x16,0x05) if ci==0 else row_bg
        add_rect(s10, cx, y, cw, 0.43, bg)
        clr = CHECK_COLOR if val=="✅" else (CROSS_COLOR if val=="❌" else WARN_COLOR)
        add_text(s10, val, cx, y+0.06, cw, 0.32, 13, color=clr, align=PP_ALIGN.CENTER)

add_text(s10, "✅ Full support  ⚠️ Partial / paid tier  ❌ Not available",
         0.4, 7.05, 12.5, 0.35, 9, italic=True, color=RGBColor(0x88,0x70,0x40), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Technical Architecture
# ════════════════════════════════════════════════════════════════════════════
s11 = blank_slide()
fill_bg(s11, DARK_BG)
add_rect(s11, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s11, 0.72)

add_text(s11, "Built Right — Technical Excellence", 0.6, 0.1, 12, 0.65, 24, bold=True, color=WHITE, font="Georgia")
add_text(s11, "A production-quality codebase built from scratch with modern tooling and thoughtful architecture.",
         0.6, 0.78, 12.5, 0.42, 12.5, color=MID_AMBER)

stack = [
    ("Frontend", "React 19 + Vite 8", "Component-based UI with fast HMR in dev and optimised production bundles."),
    ("Desktop Shell", "Electron 42", "Native Mac window chrome, IPC bridge for file I/O, single-instance lock."),
    ("Styling", "Tailwind CSS + CSS Variables", "Utility classes + 10-theme CSS variable system for instant theme switching."),
    ("Rich Text", "TipTap v3 / ProseMirror", "Full rich-text editor with font family, size, lists, and custom toolbar."),
    ("Data", "JSON + localStorage", "Zero-dependency storage — no SQLite, no ORMs, no migrations."),
    ("Bible API", "bible-api.com (cached)", "Chapter-level caching for offline operation after first load."),
    ("Build", "electron-builder + Vite", "Single DMG artefact, arm64 native for Apple Silicon."),
    ("Cross-refs", "Static JS map (31k pairs)", "Bundled at build time — zero network requests, instant lookup."),
]
for i, (layer, tech, desc) in enumerate(stack):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.4
    y = 1.38 + row * 1.38
    add_rect(s11, x, y, 6.1, 1.24, RGBColor(0x1E,0x14,0x04), RGBColor(0x6B,0x4A,0x10), 0.7)
    add_text(s11, layer, x+0.18, y+0.1, 1.5, 0.28, 9, color=MID_AMBER, bold=True)
    add_text(s11, tech, x+1.7, y+0.08, 4.2, 0.32, 11.5, color=AMBER_LIGHT, bold=True)
    add_text(s11, desc, x+0.18, y+0.52, 5.7, 0.58, 10, color=RGBColor(0xCC,0xB4,0x7A))

add_rect(s11, 0, 6.88, 13.33, 0.62, RGBColor(0x0C,0x07,0x01))
add_text(s11, "Entire app built from scratch — no template, no boilerplate, no third-party UI kit.",
         0.5, 6.95, 12.5, 0.38, 11, italic=True, color=MID_AMBER, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Vision & Roadmap
# ════════════════════════════════════════════════════════════════════════════
s12 = blank_slide()
fill_bg(s12, DARK_BG)
add_rect(s12, 0, 0, 13.33, 0.75, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s12, 0.72)

add_text(s12, "Vision & Roadmap", 0.6, 0.1, 12, 0.65, 26, bold=True, color=WHITE, font="Georgia")
add_text(s12, '"A lamp to my feet and a light to my path." — Psalm 119:105',
         0.6, 0.78, 12.5, 0.42, 13, italic=True, color=MID_AMBER)

done = [
    "✅  Full Bible reader with 31,000+ verses",
    "✅  10 biblically-named themes, all WCAG-tested",
    "✅  Cross-reference popup modal (31,000 pairs)",
    "✅  Rich-text notes editor — 9 fonts, 8 sizes",
    "✅  Verse notes linked to Bible references",
    "✅  Prayer journal with colour labels",
    "✅  Proclamations / daily declarations panel",
    "✅  iCloud Drive sync with hover timestamp",
    "✅  Emotion-based verse search",
    "✅  Highlight verses in 5 colours",
    "✅  Multiple Bible translations",
    "✅  Export notes as Word / plain text",
    "✅  Single-instance lock — no duplicate windows",
    "✅  Uniform note card heights",
]
add_rect(s12, 0.4, 1.32, 5.9, 0.44, RGBColor(0x2A,0x1A,0x06))
add_text(s12, "  ✦  Version 1.0 — Shipped", 0.55, 1.36, 5.6, 0.36, 12, bold=True, color=AMBER_LIGHT)
for i, item in enumerate(done):
    add_text(s12, item, 0.5, 1.88 + i * 0.34, 5.8, 0.32, 9.5, color=RGBColor(0xCC,0xB8,0x80))

roadmap = [
    ("📱  iOS / iPad Companion App", "Read, highlight, and pray on any device. Full iCloud sync with the Mac app."),
    ("🤖  AI Study Assistant", "Ask questions about any passage. Get commentaries, context, and cross-references via Claude AI."),
    ("📅  Reading Plans", "30-day, 90-day, and 1-year reading plans with progress tracking and daily reminders."),
    ("🌍  More Translations", "Add 20+ languages including interlinear Hebrew/Greek with Strong's concordance numbers."),
    ("👥  Shared Notes", "Share a note or prayer request with a friend or small group via AirDrop or iMessage."),
    ("🎙  Audio Bible", "Text-to-speech with multiple voices and reading speed control for accessibility."),
    ("📊  Study Insights", "Visualise your reading streaks, most-studied books, and note frequency as beautiful charts."),
]
add_rect(s12, 7.0, 1.32, 5.9, 0.44, RGBColor(0x2A,0x1A,0x06))
add_text(s12, "  ✦  Coming Next", 7.15, 1.36, 5.6, 0.36, 12, bold=True, color=AMBER_LIGHT)
for i, (title, body) in enumerate(roadmap):
    y = 1.88 + i * 0.77
    add_rect(s12, 7.0, y, 0.06, 0.48, AMBER)
    add_text(s12, title, 7.2, y, 5.6, 0.32, 11, bold=True, color=AMBER_LIGHT)
    add_text(s12, body, 7.2, y+0.3, 5.6, 0.38, 9.5, color=RGBColor(0xCC,0xB4,0x7A))

add_rect(s12, 6.56, 1.32, 0.06, 5.8, AMBER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Closing
# ════════════════════════════════════════════════════════════════════════════
s13 = blank_slide()
fill_bg(s13, DARK_BG)
add_rect(s13, 0, 0, 13.33, 2.8, RGBColor(0x1A, 0x10, 0x03))
accent_bar(s13, 0.0)

add_text(s13, "✦", 5.9, 0.4, 1.5, 1.1, 52, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_text(s13, "SEFER", 3.5, 1.15, 6.3, 1.0, 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Georgia")
add_text(s13, "Scripture. Reflection. Prayer.", 3.2, 2.0, 6.9, 0.55, 18,
         color=MID_AMBER, align=PP_ALIGN.CENTER, italic=True)

accent_bar(s13, 2.68)

add_text(s13, "\"Your word is a lamp to my feet and a light to my path.\"",
         1.5, 3.0, 10.3, 0.7, 20, italic=True, color=OFF_WHITE, align=PP_ALIGN.CENTER, font="Georgia")
add_text(s13, "— Psalm 119:105", 5.0, 3.65, 3.3, 0.45, 13,
         color=AMBER, align=PP_ALIGN.CENTER, italic=True)

add_rect(s13, 0, 4.55, 13.33, 0.04, RGBColor(0x44, 0x2E, 0x08))

summary = [
    ("31,000+\nVerses",      RGBColor(0xB0,0x7A,0x10)),
    ("10\nThemes",           RGBColor(0x6C,0x28,0xC0)),
    ("3\nModules",           RGBColor(0x1A,0x5E,0xC8)),
    ("0\nSubscriptions",     RGBColor(0x1A,0x70,0x38)),
    ("100%\nPrivate",        RGBColor(0xC0,0x18,0x18)),
]
for i, (text, color) in enumerate(summary):
    x = 0.85 + i * 2.35
    add_rect(s13, x, 4.8, 2.0, 1.6, RGBColor(0x1E,0x14,0x04), color, 1.2)
    lines = text.split('\n')
    add_text(s13, lines[0], x, 4.95, 2.0, 0.7, 26, bold=True, color=color, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s13, lines[1], x, 5.6, 2.0, 0.5, 11, color=OFF_WHITE, align=PP_ALIGN.CENTER)

add_rect(s13, 0, 6.65, 13.33, 0.85, RGBColor(0x0C,0x07,0x01))
add_text(s13, "Built with care. Designed for devotion. Yours to keep.",
         0.5, 6.72, 12.3, 0.45, 13, italic=True, color=MID_AMBER, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
out = '/Users/pradeep/Desktop/Sefer_Presentation.pptx'
prs.save(out)
print(f'Saved → {out}')
